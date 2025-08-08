
import streamlit as st
import os
import json
import re
import google.generativeai as genai
from pathlib import Path
import chromadb
from sentence_transformers import SentenceTransformer
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain.prompts import ChatPromptTemplate
from langchain_community.vectorstores import Chroma
import docx
import pptx
from llama_parse import LlamaParse
import chardet
import torch # Import torch


st.set_page_config(page_title="Document Analysis and Chat", layout="wide")


st.title("Multilingual RAG-based Document Understanding System")
st.markdown("""
    Upload a document (PDF, DOCX, TXT, PPTX) to extract key terms and definitions,
    and then chat with a RAG model to ask questions about the document's content.
""")


LLAMA_API_KEY = "" # Replace with your actual LlamaParse API key
GOOGLE_API_KEY = "" # Replace with your actual Google API key



@st.cache_resource
def load_llama_parser(api_key):
    """Loads and caches the LlamaParse parser."""
    try:
        return LlamaParse(api_key=api_key)
    except Exception as e:
        st.error(f"Error initializing LlamaParse: {e}")
        return None

def extract_text_llamaparse(file_path, parser):
    """Extracts text from PDF using LlamaParse."""
    if not parser:
        return ""
    try:
        with st.spinner(f"Parsing {os.path.basename(file_path)} with LlamaParse..."):
            parsed_docs = parser.load_data(file_path)
        full_text = "\n".join([doc.text for doc in parsed_docs])
        return full_text.strip()
    except Exception as e:
        st.error(f" Error using LlamaParse on {file_path}: {e}")
        return ""

def extract_text_from_docx(file_path):
    """Extracts text from DOCX files."""
    try:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs]).strip()
    except Exception as e:
        st.error(f" Error reading DOCX {file_path}: {e}")
        return ""

def extract_text_from_pptx(file_path):
    """Extracts text from PPTX files."""
    text = ''
    try:
        presentation = pptx.Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        st.error(f" Error reading PPTX {file_path}: {e}")
    return text.strip()

def detect_language(text):
    """Detects if the text is Arabic or English."""
    return "Arabic" if any('\u0600' <= c <= '\u06FF' for c in text) else "English"

@st.cache_data
def process_document(file_path, _llama_parser):
    """Processes an uploaded document to extract text."""
    ext = Path(file_path).suffix.lower()
    text = "" # Initialize text
    if ext == ".pdf":
        text = extract_text_llamaparse(file_path, _llama_parser)
    elif ext == ".docx":
        text = extract_text_from_docx(file_path)
    elif ext == ".pptx":
        text = extract_text_from_pptx(file_path)
    elif ext == ".txt":
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        except Exception as e:
            st.warning(f"Could not read text file with utf-8, trying chardet: {e}")
            try:
                with open(file_path, 'rb') as f:
                    raw = f.read()
                    encoding = chardet.detect(raw)['encoding']
                    text = raw.decode(encoding, errors='ignore')
            except Exception as e_chardet:
                st.error(f"Error reading text file even with chardet: {e_chardet}")
                return None
    else:
        st.warning(f" Unsupported file type: {ext}")
        return None

    if not text.strip():
        st.warning(f" No content extracted from {os.path.basename(file_path)}.")
        return None

    return {
        "file_name": os.path.basename(file_path),
        "language": detect_language(text),
        "content": text
    }

@st.cache_resource
def load_gemini_model(api_key):
    """Loads and caches the Gemini model for term extraction."""
    try:
        genai.configure(api_key=api_key)
        return genai.GenerativeModel('gemini-1.5-flash-latest')
    except Exception as e:
        st.error(f"Error initializing Gemini model: {e}")
        st.warning("Please make sure your API key is correct and you are using a valid model name.")
        return None

gemini_model = load_gemini_model(GOOGLE_API_KEY)


@st.cache_data(show_spinner="Extracting terms and definitions...")
def extract_terms_llm(content, language, _model): # Added leading underscore to model
    """Extracts terms and definitions using an LLM."""
    if not _model: # Use the renamed parameter
        return []

    prompt = f"""
    You are an expert at extracting terms and definitions from text.
    Analyze the following text and extract a list of term-definition pairs.
    A term is typically a word or phrase that is being defined, and a definition explains the meaning of the term.
    Exclude any irrelevant entries, such as URLs, single characters, photos caption, definitions with a single word  or incomplete sentences that are not definitions.
    The text is in {language}.

    Output the results as a JSON list of objects, where each object has two keys: "term" and "definition".

    Example output format:
    [
      {{"term": "Example Term", "definition": "This is an example definition."}},
      {{"term": "Another Term", "definition": "This is another definition."}}
    ]

    Text to analyze:
    ---
    {content}
    ---
    """
    try:
        response = _model.generate_content(prompt) # Use the renamed parameter
        response_text = response.text.strip()

        # Attempt to find and parse JSON within the response
        json_match = re.search(r'\[.*\]', response_text, re.DOTALL)
        if json_match:
            json_string = json_match.group(0)
            try:
                extracted_data = json.loads(json_string)
                # Basic filtering for quality
                filtered_data = [
                    item for item in extracted_data
                    if item.get('term') and len(item['term'].strip()) > 1
                    and item.get('definition') and len(item['definition'].strip().split()) > 1 # Definition has more than one word
                ]
                return filtered_data
            except json.JSONDecodeError as e:
                st.error(f" Could not decode JSON from LLM response: {e}")
                st.text(f"Raw LLM response: {response_text}")
                return []
        else:
             st.warning(" Could not find JSON in LLM response.")
             st.text(f"Raw LLM response: {response_text}")
             return []

    except Exception as e:
        st.error(f" Error during term extraction with LLM: {e}")
        return []

@st.cache_resource
def load_embedding_model():
    """Loads and caches the Sentence Transformer embedding model."""
    try:
        return SentenceTransformer("paraphrase-multilingual-mpnet-base-v2")
    except Exception as e:
        st.error(f"Error initializing Sentence Transformer model: {e}")
        return None

embedding_model = load_embedding_model()


class SentenceTransformerEmbeddingWrapper:
     def __init__(self, model):
        self.model = model
        self.name = "SentenceTransformerEmbeddingWrapper" # Added name attribute


     def embed_documents(self, texts, batch_size=32): # Added batch_size parameter
        """Embeds a list of documents in batches."""
        if isinstance(texts, str): # Ensure texts is a list
            texts = [texts]
        # Add a check for empty list
        if not texts:
            return []

        embeddings = []
        try:
            for i in range(0, len(texts), batch_size):
                batch = texts[i:i+batch_size]
                batch_embeddings = self.model.encode(batch).tolist()
                embeddings.extend(batch_embeddings)
                # Clear GPU memory cache after each batch
                if torch.cuda.is_available():
                    torch.cuda.empty_cache()
            return embeddings
        except Exception as e:
            st.error(f"Error during document embedding: {e}")
            return []

     def embed_query(self, text):
        """Embeds a single query."""
        try:
            return self.model.encode([text])[0].tolist()
        except Exception as e:
            st.error(f"Error during query embedding: {e}")
            return []

sentence_transformer_embedding_function = None
if embedding_model:
    sentence_transformer_embedding_function = SentenceTransformerEmbeddingWrapper(embedding_model)

@st.cache_resource
def load_chat_llm(api_key):
     """Loads and caches the Chat LLM."""
     try:
        return ChatGoogleGenerativeAI(model="gemini-1.5-flash-latest", google_api_key=api_key)
     except Exception as e:
        st.error(f"Error initializing Chat LLM: {e}")
        return None

chat_llm = load_chat_llm(GOOGLE_API_KEY)





uploaded_file = st.file_uploader("Choose a file", type=["pdf", "docx", "txt", "pptx"])

temp_file_path = None
document_data = None # Store processed document data
collection_name = "raw_file_chunks" # Define collection name


@st.cache_data
def process_and_store_document(uploaded_file, _llama_parser, _embedding_func, _llm_chat):
    """Processes uploaded file, extracts text, and stores chunks in vector db."""
    temp_dir = "temp_uploads"
    os.makedirs(temp_dir, exist_ok=True)
    temp_file_path = os.path.join(temp_dir, uploaded_file.name)
    with open(temp_file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    st.success(f"File '{uploaded_file.name}' uploaded successfully!")

    document_data = process_document(temp_file_path, _llama_parser)

    if document_data and document_data.get("content"):
        with st.spinner("Storing document chunks in vector database..."):
            try:
                chroma_client = chromadb.PersistentClient(path="vector_db")
                # Ensure collection name is unique per file or handle updates
                file_collection_name = f"chunks_{re.sub('[^a-zA-Z0-9_]', '', uploaded_file.name)}"

                text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=500)
                chunks = text_splitter.split_text(document_data["content"])
                texts, metadatas, ids = [], [], []
                uid = 0
                for i, chunk in enumerate(chunks):
                    texts.append(chunk)
                    metadatas.append({"file": document_data['file_name'], "chunk_number": i + 1})
                    ids.append(f"raw-{document_data['file_name']}-{uid}")
                    uid += 1

                if texts and _embedding_func:
                    embeddings = _embedding_func.embed_documents(texts)
                    if embeddings:
                        # Initialize collection WITHOUT embedding_function here
                        collection = chroma_client.get_or_create_collection(name=file_collection_name)

                        def batch_add(collection, docs, embs, metas, ids, batch_size=1000):
                            for i in range(0, len(docs), batch_size):
                                collection.add(
                                    documents=docs[i:i+batch_size],
                                    embeddings=embs[i:i+batch_size],
                                    metadatas=metas[i:i+batch_size],
                                    ids=ids[i:i+batch_size]
                                )
                              

                        batch_add(collection, texts, embeddings, metadatas, ids)
                        st.success(f" Stored {len(texts)} chunks from {document_data['file_name']} in vector database collection: {file_collection_name}.")

                 
                        return document_data, file_collection_name
                    else:
                        st.error("Could not generate embeddings for document chunks.")
                        return document_data, None
                else:
                     st.warning(" No chunks generated for vectorization or embedding function not available.")
                     return document_data, None
            except Exception as e:
                st.error(f"Error during embedding or storing raw file chunks: {e}")
                return document_data, None
    else:
        return document_data, None



if uploaded_file is not None:
    llama_parser = load_llama_parser(LLAMA_API_KEY)
    if llama_parser and embedding_model and chat_llm:
         processed_data, stored_collection_name = process_and_store_document(uploaded_file, llama_parser, sentence_transformer_embedding_function, chat_llm)
         if processed_data:
             st.session_state['document_data'] = processed_data
             st.session_state['stored_collection_name'] = stored_collection_name
             # Initialize chat history if not already present
             if 'chat_history' not in st.session_state:
                 st.session_state['chat_history'] = []



st.header("Extracted Terms and Definitions")
if 'document_data' in st.session_state and st.session_state['document_data'] and st.session_state['document_data'].get("content"):
    document_data_display = st.session_state['document_data']
    st.write(f"Processing terms and definitions for: {document_data_display['file_name']}")


    extracted_terms = extract_terms_llm(
        document_data_display["content"],
        document_data_display["language"],
        gemini_model
    )

    if extracted_terms:
        st.subheader("Terms and Definitions:")
    
        for item in extracted_terms:
            term = item.get("term", "N/A")
            definition = item.get("definition", "N/A")
            if term != "N/A" and definition != "N/A":
                 st.markdown(f"- **{term}**: {definition}")
    else:
        st.info("No terms and definitions extracted.")
else:
    st.write("Upload a document to see extracted terms and definitions here.")

st.header("Chat with the Document")
if 'document_data' in st.session_state and st.session_state['document_data'] and st.session_state['document_data'].get("content") and 'stored_collection_name' in st.session_state and st.session_state['stored_collection_name'] and chat_llm and embedding_model:

    # Initialize or get RAG chain from session state to avoid re-creating
    if 'rag_chain' not in st.session_state or st.session_state['rag_chain'] is None:
        try:
             chroma_client = chromadb.PersistentClient(path="vector_db")
             vectorstore = Chroma(
                collection_name=st.session_state['stored_collection_name'],
                embedding_function=sentence_transformer_embedding_function, # Use the cached embedding function
                persist_directory="vector_db"
             )
             retriever = vectorstore.as_retriever(search_kwargs={"k": 5})

             prompt_template = ChatPromptTemplate.from_template("""
             You are an expert technical assistant.
             Answer the question clearly and completely, using ONLY the context below.
             If the user inputs a misspelled question, try to correct it by yourself looking at the closest correct spelling that exists in context and answer accordingly.
             if you can not find context in the text answer and say that the texts do not include an answer. If the context in Arabic and the question is in Arabic answer in Arabic according to context which is in Arabic”

             Include:
             - Definitions if relevant
             - Examples or units of measurement
             - Any formulas mentioned in the context
             - Only include the most relevant information without any extra output other than information that exist in the context


             Context:
             {context}

             Question: {input}
             Answer:
             """)
             document_chain = create_stuff_documents_chain(llm=chat_llm, prompt=prompt_template)
             st.session_state['rag_chain'] = create_retrieval_chain(retriever=retriever, combine_docs_chain=document_chain)
             st.success("RAG chain initialized. You can now chat with the document.")
          
             st.session_state['retriever'] = retriever
        except Exception as e:
            st.error(f"Error setting up RAG chain: {e}")
            st.session_state['rag_chain'] = None # Ensure rag_chain is None on failure
            st.session_state['retriever'] = None


    # Display chat history (restoring this)
    if 'chat_history' in st.session_state:
        for chat_entry in st.session_state['chat_history']:
            st.markdown(f"**You:** {chat_entry['question']}")
            st.markdown(f"**AI:** {chat_entry['answer']}")


    user_question = st.text_input("Ask a question about the document:", key="chat_input") # Added a key

    if st.button("Ask", key="ask_button"): # Added a key
        if user_question:
            if 'rag_chain' in st.session_state and st.session_state['rag_chain']:
                with st.spinner("Getting answer from the document..."):
                    try:
                        result = st.session_state['rag_chain'].invoke({"input": user_question})

                    
                        answer = result.get("answer", "Could not retrieve an answer.")

                        # Append to chat history (restoring this)
                        if 'chat_history' not in st.session_state:
                             st.session_state['chat_history'] = []
                        st.session_state['chat_history'].append({"question": user_question, "answer": answer})


                        st.rerun() # Keep rerun to clear input and update page with chat history

                    except Exception as e:
                        st.error(f"Error during RAG chain invocation: {e}")
            else:
                 st.warning("RAG chain is not initialized. Please upload a document and ensure all models loaded successfully.")
        else:
            st.warning("Please enter a question.")

else:
    if 'document_data' not in st.session_state or not st.session_state['document_data'] or not st.session_state['document_data'].get("content"):
        st.write("Upload a document to enable the chat interface.")
    elif 'stored_collection_name' not in st.session_state or not st.session_state['stored_collection_name']:
         st.warning("Document processed, but storage in vector database failed. Chat disabled.")
    elif not chat_llm or not embedding_model:

         st.warning("LLM or embedding model not initialized. Chat disabled.")
