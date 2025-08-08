from pathlib import Path
import pytesseract
from PIL import Image
import docx
import json
import os
import pptx 
from llama_parse import LlamaParse

LLAMA_API_KEY = "llx-QxOOQzstU9I7CYUBpEuzw8jDrJCh3FSDIcEZZjWSwhWWEAL5"  

def extract_text_llamaparse(file_path):
    parser = LlamaParse(api_key=LLAMA_API_KEY)
    try:
        parsed_docs = parser.load_data(file_path)
        full_text = "\n".join([doc.text for doc in parsed_docs])
        return full_text.strip()
    except Exception as e:
        print(f" Error using LlamaParse on {file_path}: {e}")
        return ""


def extract_text_from_docx(file_path):
    try:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs]).strip()
    except Exception as e:
        print(f" Error reading DOCX {file_path}: {e}")
        return ""


def extract_text_from_pptx(file_path):
    text = ''
    try:
        presentation = pptx.Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f" Error reading PPTX {file_path}: {e}")
    return text.strip()


def detect_language(text):
    return "Arabic" if any('\u0600' <= c <= '\u06FF' for c in text) else "English"

def process_document(text, file_name, language):
    if not text.strip():
        print(f"❌ No content in uploaded file: {file_name}")
        return None

    return {
        "file_name": file_name,
        "language": language,
        "content": text.strip()
    }


def process_all_documents(folder_path, output_json="output.json"):
    results = []
    
    for file in os.listdir(folder_path):
        if file.lower().endswith(('.pdf', '.docx', '.txt', '.pptx')):
            doc_data = process_document(os.path.join(folder_path, file))
            if doc_data:
                results.append(doc_data)

    with open(output_json, "w", encoding="utf-8") as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    print(f" Done. Saved to {output_json}")
process_all_documents(r"D:\ragg")
import json
import re

output_json_path = "output.json"

try:
    with open(output_json_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    print(f" Successfully loaded data from {output_json_path}")
except (FileNotFoundError, json.JSONDecodeError) as e:
    print(f" Error loading data from {output_json_path}: {e}")
    data = [] 

# Iterate and clean document content
#cleaned_data = []
#for document in data:
    #content = document.get("content", "") # Get content, default to empty string

    # Apply cleaning steps
    # Remove leading/trailing whitespace from each line
    #cleaned_content_lines = [line.strip() for line in content.splitlines()]
    # Join lines, removing extra blank lines but preserving paragraph breaks (heuristic)
    # Keep single blank lines, remove multiple blank lines
    #cleaned_content = "\n".join(filter(lambda x: x or (x == "" and cleaned_content_lines[cleaned_content_lines.index(x)+1] != ""), cleaned_content_lines))

    # Replace multiple spaces with a single space
    #cleaned_content = re.sub(r'\s+', ' ', cleaned_content).strip()

    # Update the document with the cleaned content
   # document["content"] = cleaned_content
   # cleaned_data.append(document)

#print(" Cleaning of document content complete.")
#cleaned_output_json_path = "cleaned_data.json"

#فry:
    #with open(cleaned_output_json_path, "w", encoding="utf-8") as f:
       # json.dump(cleaned_data, f, ensure_ascii=False, indent=2)
    #print(f" Successfully saved cleaned data to {cleaned_output_json_path}")
#except Exception as e:
    #print(f" Error saving cleaned data to {cleaned_output_json_path}: {e}")
                
#output_json_path = "output.json"

with open(output_json_path, "r", encoding="utf-8") as f:
    data = json.load(f)



import google.generativeai as genai
import json


GOOGLE_API_KEY = "AIzaSyBGZpCRyC548KpqkHf5qYlICO0EfSz7EcY"


try:
    genai.configure(api_key=GOOGLE_API_KEY)
    gemini_model = genai.GenerativeModel('gemini-1.5-flash-latest')
except Exception as e:
    print(f"Error initializing Gemini model: {e}")
    print("Please make sure your API key is correct and you are using a valid model name.")
    gemini_model = None 


all_extracted_terms_llm = []

if gemini_model:
    for document in data:
        content = document["content"]
        file_name = document["file_name"]
        language = document["language"]

        print(f"Extracting terms and definitions from: {file_name} ({language})")

        
        prompt = f"""
        You are an expert at extracting terms and definitions from text.
        Analyze the following text and extract a list of term-definition pairs.
        A term is typically a word or phrase that is being defined, and a definition explains the meaning of the term.
        Exclude any irrelevant entries, such as URLs, single characters, photos caption, definitions with a single word  or incomplete sentences that are not definitions.
        The text is in {language}.

        Output the results as a JSON list of objects, where each object has two keys: "term" and "definition".

        Example output format:
        [
          {{"term": "Example Term", "definition": "This is an example definition."}},
          {{"term": "Another Term", "definition": "This is another definition."}}
        ]

        Text to analyze:
        ---
        {content}
        ---
        """

        try:
           
            response = gemini_model.generate_content(prompt)

            
            response_text = response.text.strip()
            json_start = response_text.find('[')
            json_end = response_text.rfind(']')

            extracted_data = []
            if json_start != -1 and json_end != -1:
                json_string = response_text[json_start : json_end + 1]
                try:
                    extracted_data = json.loads(json_string)
                except json.JSONDecodeError as e:
                    print(f" Could not decode JSON from LLM response for {file_name}: {e}")
                    print(f"Raw LLM response: {response_text}")
            else:
                 print(f" Could not find JSON in LLM response for {file_name}.")
                 print(f"Raw LLM response: {response_text}")


            all_extracted_terms_llm.append({
                "file_name": file_name,
                "extracted_terms": extracted_data
            })

        except Exception as e:
            print(f" Error processing {file_name} with LLM: {e}")

    print("\n LLM extraction complete for all documents.")
else:
    print("\nSkipping LLM extraction due to model initialization failure.")

print(prompt)

consolidated_terms_llm = {}

for document_data in all_extracted_terms_llm:
    file_name = document_data['file_name']
    extracted_terms = document_data.get('extracted_terms', []) 

    for entry in extracted_terms:
        
        term = entry.get('term', '').lower().strip()  
        definition = entry.get('definition', '').strip()

        if term and definition:
            if term not in consolidated_terms_llm:
                consolidated_terms_llm[term] = []

         
            if definition not in consolidated_terms_llm[term]:
                consolidated_terms_llm[term].append(definition)

print(" Consolidation of LLM extracted terms complete.")


for term, definitions in consolidated_terms_llm.items():
     print(f"Term: {term}")
     for i, definition in enumerate(definitions):
         print(f"  Definition {i+1}: {definition}")


terms_with_multiple_definitions_llm_by_language = {}

for document_data in all_extracted_terms_llm:
    file_name = document_data['file_name']
    extracted_terms = document_data.get('extracted_terms', [])
    language = document_data.get('language', 'unknown') 

    if language not in terms_with_multiple_definitions_llm_by_language:
        terms_with_multiple_definitions_llm_by_language[language] = {}

    for entry in extracted_terms:
        term = entry.get('term', '').lower().strip()
        definition = entry.get('definition', '').strip()

        if term and definition:
            if term not in terms_with_multiple_definitions_llm_by_language[language]:
                terms_with_multiple_definitions_llm_by_language[language][term] = []

            if definition not in terms_with_multiple_definitions_llm_by_language[language][term]:
                terms_with_multiple_definitions_llm_by_language[language][term].append(definition)


terms_with_multiple_definitions_llm_by_language_filtered = {}
for language, terms_data in terms_with_multiple_definitions_llm_by_language.items():
    terms_with_multiple_definitions_llm_by_language_filtered[language] = {
        term: definitions for term, definitions in terms_data.items() if len(definitions) > 1
    }


print(" Identification of terms with multiple definitions by language complete.")



print("--- Terms with different definitions (LLM) by Language ---")
found_multiple_definitions = False
for language, terms_data in terms_with_multiple_definitions_llm_by_language_filtered.items():
    if terms_data:
        found_multiple_definitions = True
        print(f"\n--- Language: {language} ---")
        for term, definitions in terms_data.items():
            print(f"\nTerm: {term}")
            for i, definition in enumerate(definitions):
                print(f"  Definition {i+1}: {definition}")

if not found_multiple_definitions:
    print("\nNo terms with multiple definitions found within the same language using the LLM extraction.")



import arabic_reshaper
from bidi.algorithm import get_display

for term, definitions in consolidated_terms_llm.items():
    reshaped_term = arabic_reshaper.reshape(f"Term: {term}")
    print(get_display(reshaped_term))
    
    if len(definitions) > 1:
        print("  Note: This term has differing definitions:")
        for i, definition in enumerate(definitions):
            reshaped_def = arabic_reshaper.reshape(f"    Definition {i+1}: {definition}")
            print(get_display(reshaped_def))
    else:
        if definitions:
            reshaped_def = arabic_reshaper.reshape(f"  Definition: {definitions[0]}")
            print(get_display(reshaped_def))
        else:
            reshaped_def = arabic_reshaper.reshape("  Definition: [No definition found]")
            print(get_display(reshaped_def))\


os.makedirs("outputs", exist_ok=True)

with open("outputs/generated_definitions.json", "w", encoding="utf-8") as f:
    json.dump(all_extracted_terms_llm, f, ensure_ascii=False, indent=2)
print(" All extracted terms and definitions saved to outputs/generated_definitions.json")
import os
import json
import arabic_reshaper
from bidi.algorithm import get_display


terms_with_differences = {}

for term, definitions in consolidated_terms_llm.items():
    reshaped_term = arabic_reshaper.reshape(f"Term: {term}")
    print(get_display(reshaped_term))

    if len(definitions) > 1:
        print("  Note: This term has differing definitions:")
        terms_with_differences[term] = definitions  
        for i, definition in enumerate(definitions):
            reshaped_def = arabic_reshaper.reshape(f"    Definition {i+1}: {definition}")
            print(get_display(reshaped_def))
    else:
        if definitions:
            reshaped_def = arabic_reshaper.reshape(f"  Definition: {definitions[0]}")
            print(get_display(reshaped_def))
        else:
            reshaped_def = arabic_reshaper.reshape("  Definition: [No definition found]")
            print(get_display(reshaped_def))


os.makedirs("outputs", exist_ok=True)


with open("outputs/generated_definitions.json", "w", encoding="utf-8") as f:
    json.dump(all_extracted_terms_llm, f, ensure_ascii=False, indent=2)


with open("outputs/differing_definitions.json", "w", encoding="utf-8") as f:
    json.dump(terms_with_differences, f, ensure_ascii=False, indent=2)


