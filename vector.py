
import os
import json
from sentence_transformers import SentenceTransformer
from langchain.text_splitter import RecursiveCharacterTextSplitter
import chromadb


embedding_model = SentenceTransformer("sentence-transformers/paraphrase-multilingual-mpnet-base-v2")
chroma_client = chromadb.PersistentClient(path="vector_db")
text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=500)


def extract_text_from_file(file_path):
    ext = file_path.lower()
    try:
        if ext.endswith(".txt"):
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()

        elif ext.endswith(".pdf"):
            from llama_parse import LlamaParse
           
            parser = LlamaParse(api_key="")
            parsed_docs = parser.load_data(file_path)
            return "\n".join([doc.text for doc in parsed_docs])

        elif ext.endswith(".docx"):
            from docx import Document
            doc = Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs])

        elif ext.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)

        else:
            return ""
    except Exception as e:
        print(f" Error reading {file_path}: {e}")
        return ""


def store_glossary(file_path, collection_name):
    with open(file_path, 'r', encoding='utf-8') as f:
        glossary_data = json.load(f)

    texts, metadatas, ids = [], [], []
    uid = 0

    for entry in glossary_data:
        file_name = entry['file_name']
        for item in entry.get('extracted_terms', []):
            term = item.get('term', '').strip()
            definition = item.get('definition', '').strip()
            if term and definition:
                texts.append(f"{term}: {definition}")
                metadatas.append({"file": file_name, "term": term})
                ids.append(f"term-{uid}")
                uid += 1

    if texts:
        embeddings = embedding_model.encode(texts).tolist()
        collection = chroma_client.get_or_create_collection(name=collection_name)
        collection.add(documents=texts, embeddings=embeddings, metadatas=metadatas, ids=ids)
        print(f" Stored: {collection_name} ({len(texts)} items)")
    else:
        print(f" No valid terms found in {file_path}")


def store_conflicting_terms(file_path, collection_name):
    with open(file_path, 'r', encoding='utf-8') as f:
        differing_data = json.load(f)

    texts, metadatas, ids = [], [], []
    uid = 0

    for term, definitions in differing_data.items():
        for i, definition in enumerate(definitions):
            texts.append(f"{term}: {definition}")
            metadatas.append({"term": term, "version": i + 1})
            ids.append(f"diff-{uid}")
            uid += 1

    if texts:
        embeddings = embedding_model.encode(texts).tolist()
        collection = chroma_client.get_or_create_collection(name=collection_name)
        collection.add(documents=texts, embeddings=embeddings, metadatas=metadatas, ids=ids)
        print(f" Stored: {collection_name} ({len(texts)} items)")
    else:
        print(f" No conflicting definitions found in {file_path}")


def store_raw_file_chunks(data_dir="data", collection_name="raw_file_chunks"):
    texts, metadatas, ids = [], [], []
    uid = 0

    for filename in os.listdir(data_dir):
        file_path = os.path.join(data_dir, filename)
        if filename.lower().endswith(('.pdf', '.docx', '.pptx', '.txt')):
            full_text = extract_text_from_file(file_path).strip()
            if full_text:
                chunks = text_splitter.split_text(full_text)
                for i, chunk in enumerate(chunks):
                    texts.append(chunk)
                    metadatas.append({"file": filename, "chunk_number": i + 1})
                    ids.append(f"raw-{uid}")
                    uid += 1

    if texts:
        embeddings = embedding_model.encode(texts).tolist()
        collection = chroma_client.get_or_create_collection(name=collection_name)

        def batch_add(collection, docs, embs, metas, ids, batch_size=5000):
            for i in range(0, len(docs), batch_size):
                collection.add(
                    documents=docs[i:i+batch_size],
                    embeddings=embs[i:i+batch_size],
                    metadatas=metas[i:i+batch_size],
                    ids=ids[i:i+batch_size]
                )
                print(f" Added batch {i // batch_size + 1}")

        batch_add(collection, texts, embeddings, metadatas, ids)
        print(f" Stored: {collection_name} ({len(texts)} chunks)")
    else:
        print(" No raw content found for vectorization.")


if __name__ == "__main__":
    
    store_glossary("outputs/generated_definitions.json", "terms_and_definitions")
    store_conflicting_terms("outputs/differing_definitions.json", "conflicting_terms")


    store_raw_file_chunks("data", "raw_file_chunks")
